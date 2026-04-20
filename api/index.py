import os
import base64
from pathlib import Path
from io import BytesIO
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from typing import Optional, List
from dotenv import load_dotenv

from api.database import (
    init_db, create_session, get_sessions, get_session, 
    add_message, update_session_title, delete_session,
    save_lecture_note, get_lecture_notes, delete_lecture_note, get_lecture_note_by_id
)

# Load .env from project root (works locally and on production)
env_path = Path(__file__).parent.parent / ".env"
load_dotenv(env_path)

GEMINI_KEY = os.getenv("GEMINI_API_KEY", "")
os.environ["GEMINI_API_KEY"] = GEMINI_KEY

app = FastAPI()

init_db()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Import both services
from api.services.llm_service import process_multimodal_query
from api.services.langchain_service import (
    process_rag_query, 
    generate_quiz_with_rag,
    generate_flashcards_with_rag,
    get_rag_system
)

GEMINI_KEY = os.getenv("GEMINI_API_KEY", "")
print(f"Index.py - Env path: {env_path}")
print(f"Index.py - GEMINI_KEY: {GEMINI_KEY[:20] if GEMINI_KEY else 'NOT FOUND'}")

os.environ["GEMINI_API_KEY"] = GEMINI_KEY

MAX_FILE_SIZE = 50 * 1024 * 1024

@app.post("/api/chat")
async def chat_endpoint(
    message: str = Form(...),
    user: str = Form(default="User"),
    session_id: Optional[int] = Form(default=None),
    save_history: bool = Form(default=True),
    files: Optional[List[UploadFile]] = File(default=None)
):
    try:
        processed_files = []
        
        if files:
            for file in files:
                content = await file.read()
                
                if len(content) > MAX_FILE_SIZE:
                    return JSONResponse({
                        "error": f"File '{file.filename}' exceeds 50MB limit. Please use smaller files."
                    }, status_code=413)
                
                file_type = file.content_type or "application/octet-stream"
                
                if file_type.startswith("image/"):
                    b64_data = base64.b64encode(content).decode("utf-8")
                    processed_files.append({
                        "type": "image",
                        "name": file.filename,
                        "data": b64_data,
                        "mime": file_type
                    })
                elif file_type.startswith("audio/"):
                    b64_data = base64.b64encode(content).decode("utf-8")
                    processed_files.append({
                        "type": "audio",
                        "name": file.filename,
                        "data": b64_data,
                        "mime": file_type
                    })
                elif file_type.startswith("video/"):
                    return JSONResponse({
                        "error": "Video processing requires external storage. For now, please provide a video URL or extract audio."
                    }, status_code=400)
                elif file.filename.endswith(".pdf"):
                    text_content = ""
                    extraction_method = "none"
                    
                    try:
                        from pypdf import PdfReader
                        from io import BytesIO
                        pdf_file = BytesIO(content)
                        reader = PdfReader(pdf_file)
                        for page in reader.pages:
                            text_content += page.extract_text() or ""
                        if text_content.strip():
                            extraction_method = "pypdf"
                    except Exception as e:
                        print(f"PDF pypdf error: {e}")
                    
                    if not text_content.strip():
                        try:
                            import pdfplumber
                            from io import BytesIO
                            with pdfplumber.open(BytesIO(content)) as pdf:
                                for page in pdf.pages:
                                    txt = page.extract_text()
                                    if txt:
                                        text_content += txt + "\n"
                            if text_content.strip():
                                extraction_method = "pdfplumber"
                        except Exception as e:
                            print(f"PDF pdfplumber error: {e}")
                    
                    print(f"=== PDF EXTRACT: {file.filename} using {extraction_method}, got {len(text_content)} chars ===")
                    print(f"=== PDF SAMPLE: {repr(text_content[:300])} ===")
                    
                    if text_content.strip() and len(text_content) > 30:
                        processed_files.append({
                            "type": "text",
                            "name": file.filename,
                            "data": text_content
                        })
                    else:
                        print(f"=== PDF FAILED: setting placeholder ===")
                        processed_files.append({
                            "type": "text",
                            "name": file.filename,
                            "data": f"[PDF file '{file.filename}' uploaded but text extraction failed - please provide the file as text/doc or use a different PDF]"
                        })
                elif file.filename.endswith((".txt", ".md")):
                    text_content = content.decode("utf-8", errors="ignore")
                    processed_files.append({
                        "type": "text",
                        "name": file.filename,
                        "data": text_content
                    })
                elif file.filename.endswith((".doc", ".docx")):
                    try:
                        from docx import Document
                        from io import BytesIO
                        doc_file = BytesIO(content)
                        doc = Document(doc_file)
                        text_content = "\n".join([p.text for p in doc.paragraphs])
                        processed_files.append({
                            "type": "text",
                            "name": file.filename,
                            "data": text_content
                        })
                    except Exception as e:
                        print(f"Error extracting docx text: {e}")
                        text_content = f"[Could not extract text from {file.filename}]"
                        processed_files.append({
                            "type": "text",
                            "name": file.filename,
                            "data": text_content
                        })
                elif file.filename.endswith((".ppt", ".pptx")):
                    try:
                        from pptx import Presentation
                        from io import BytesIO
                        ppt_file = BytesIO(content)
                        prs = Presentation(ppt_file)
                        text_content = ""
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    for para in shape.text_frame.paragraphs:
                                        text_content += para.text + "\n"
                        if text_content.strip():
                            processed_files.append({
                                "type": "text",
                                "name": file.filename,
                                "data": text_content
                            })
                        else:
                            processed_files.append({
                                "type": "text",
                                "name": file.filename,
                                "data": "[PowerPoint is empty]"
                            })
                    except Exception as e:
                        print(f"Error extracting pptx text: {e}")
                        text_content = f"[Could not extract text from {file.filename}]"
                        processed_files.append({
                            "type": "text",
                            "name": file.filename,
                            "data": text_content
                        })
                else:
                    processed_files.append({
                        "type": "text",
                        "name": file.filename,
                        "data": content.decode("utf-8", errors="ignore")[:1000]
                    })
        
        # Use LangChain RAG if files attached, otherwise use regular LLM
        if processed_files:
            # RAG mode with document context
            try:
                result = await process_rag_query(message, user, processed_files, session_id or 0)
            except Exception as e:
                print(f"RAG error, falling back to regular LLM: {e}")
                result = await process_multimodal_query(message, user, processed_files)
        else:
            # Regular LLM mode
            result = await process_multimodal_query(message, user, processed_files)
        
        if save_history:
            if not session_id:
                session_id = create_session(user, message[:30] + "...")
            add_message(session_id, "user", message)
            add_message(session_id, "bot", result)
        
        return {"reply": result, "session_id": session_id}
        
    except Exception as e:
        return JSONResponse({
            "error": f"Error processing request: {str(e)}"
        }, status_code=500)

@app.get("/api/health")
async def health_check():
    return {"status": "healthy", "service": "EduChat API"}

@app.get("/api/chat/history")
async def get_chat_history(user: str = "User"):
    try:
        sessions = get_sessions(user)
        return {"sessions": sessions}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/api/chat/session/{session_id}")
async def get_chat_session(session_id: int):
    try:
        session = get_session(session_id)
        if not session:
            return JSONResponse({"error": "Session not found"}, status_code=404)
        return session
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/chat/session")
async def create_chat_session(user: str = Form(...), title: str = Form(default="New Chat")):
    try:
        print(f"Creating session for user: {user}, title: {title}")
        session_id = create_session(user, title)
        print(f"Session created with ID: {session_id}")
        return {"session_id": session_id, "title": title}
    except Exception as e:
        print(f"Error creating session: {e}")
        return JSONResponse({"error": str(e)}, status_code=500)

@app.put("/api/chat/session/{session_id}")
async def update_chat_session(session_id: int, title: str = Form(...)):
    try:
        update_session_title(session_id, title)
        return {"success": True}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.delete("/api/chat/session/{session_id}")
async def delete_chat_session(session_id: int):
    try:
        delete_session(session_id)
        return {"success": True}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

def clean_extracted_text(text: str) -> str:
    import re
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    text = re.sub(r'([a-z])\n([a-z])', r'\1 \2', text, flags=re.IGNORECASE)
    text = re.sub(r' {2,}', ' ', text)
    text = re.sub(r'\n ', '\n', text)
    text = re.sub(r' \n', '\n', text)
    text = re.sub(r'(\w)-\s+(\w)', r'\1\2', text)
    text = text.strip()
    return text

@app.post("/api/extract-pdf")
async def extract_pdf(content: str = Form(...), file_type: str = Form(default="pdf")):
    try:
        import base64
        from io import BytesIO
        from fastapi.responses import StreamingResponse
        import json
        import asyncio
        
        if file_type == "pdf":
            from pypdf import PdfReader
            
            async def generate_pages():
                pdf_bytes = base64.b64decode(content)
                pdf_file = BytesIO(pdf_bytes)
                reader = PdfReader(pdf_file)
                total_pages = len(reader.pages)
                
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    text = clean_extracted_text(text)
                    
                    page_data = {
                        "page_num": i + 1,
                        "total_pages": total_pages,
                        "text": text,
                        "extracted": True
                    }
                    yield f"data: {json.dumps(page_data)}\n\n"
                    await asyncio.sleep(0.1)
                
                yield f"data: {json.dumps({'done': True})}\n\n"
            
            return StreamingResponse(generate_pages(), media_type="text/event-stream")
        
        elif file_type in ["doc", "docx"]:
            from docx import Document
            doc_bytes = base64.b64decode(content)
            doc_file = BytesIO(doc_bytes)
            doc = Document(doc_file)
            
            text = "\n".join([para.text for para in doc.paragraphs])
            text = clean_extracted_text(text)
            return {"pages": [{"page_num": 1, "text": text}], "total_pages": 1, "success": True}
        
        elif file_type == "txt" or file_type == "md":
            text = base64.b64decode(content).decode('utf-8', errors='ignore')
            text = clean_extracted_text(text)
            return {"pages": [{"page_num": 1, "text": text}], "total_pages": 1, "success": True}
        
        elif file_type in ["ppt", "pptx"]:
            try:
                from pptx import Presentation
                doc_bytes = base64.b64decode(content)
                doc_file = BytesIO(doc_bytes)
                prs = Presentation(doc_file)
                
                text = ""
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = f"\n=== Slide {slide_num} ===\n\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            for para in shape.text_frame.paragraphs:
                                para_text = para.text.strip()
                                if para_text:
                                    slide_text += para_text + "\n"
                        elif hasattr(shape, "text") and shape.text:
                            slide_text += shape.text.strip() + "\n"
                    if slide_text.strip():
                        text += slide_text + "\n"
                
                text = clean_extracted_text(text)
                return {"pages": [{"page_num": 1, "text": text}], "total_pages": len(prs.slides), "success": True}
            except Exception as e:
                return {"error": f"Failed to extract PowerPoint: {str(e)}"}
        
        else:
            return {"error": "Unsupported format"}
            
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/api/notes")
async def get_notes(user: str = "User"):
    try:
        notes = get_lecture_notes(user)
        return {"notes": notes}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/api/notes/{note_id}/content")
async def get_note_content(note_id: int):
    try:
        from api.database import get_lecture_note_by_id
        note = get_lecture_note_by_id(note_id)
        if not note:
            return JSONResponse({"error": "Note not found"}, status_code=404)
        
        import base64
        from io import BytesIO
        from fastapi.responses import StreamingResponse
        import json
        import asyncio
        
        content = note.get("content", "")
        file_type = note.get("file_type", "")
        name = note.get("name", "")
        
        print(f"Note ID: {note_id}, FileType: {file_type}, Content length: {len(content)}, Name: {name}")
        
        def try_decode_base64(s):
            if not s:
                return None
            try:
                decoded = base64.b64decode(s, validate=True)
                return decoded
            except:
                return None
        
        decoded_content = try_decode_base64(content)
        is_encoded = decoded_content is not None
        print(f"Is encoded: {is_encoded}, Content preview: {content[:50] if content else 'empty'}")
        
        async def generate_extraction():
            try:
                # For non-encoded content (plain text), return directly
                if not is_encoded:
                    yield f"data: {json.dumps({'page_num': 1, 'total_pages': 1, 'text': content})}\n\n"
                    yield f"data: {json.dumps({'done': True, 'note_name': name})}\n\n"
                    return
                
                if file_type == "pdf":
                    from pypdf import PdfReader
                    pdf_bytes = decoded_content if is_encoded else content.encode()
                    pdf_file = BytesIO(pdf_bytes)
                    reader = PdfReader(pdf_file)
                    total_pages = len(reader.pages)
                    
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text() or ""
                        text = clean_extracted_text(text)
                        page_text = f"\n\n---\n**Page {i + 1}**\n---\n\n{text}"
                        
                        page_data = {
                            "page_num": i + 1,
                            "total_pages": total_pages,
                            "text": page_text,
                            "extracted": True
                        }
                        yield f"data: {json.dumps(page_data)}\n\n"
                        await asyncio.sleep(0.05)
                    
                    yield f"data: {json.dumps({'done': True})}\n\n"
                
                elif file_type in ["doc", "docx"]:
                    from docx import Document
                    doc_bytes = decoded_content if is_encoded else content.encode()
                    doc_file = BytesIO(doc_bytes)
                    doc = Document(doc_file)
                    
                    text = "\n".join([para.text for para in doc.paragraphs])
                    text = clean_extracted_text(text)
                    
                    page_data = {"page_num": 1, "total_pages": 1, "text": text, "extracted": True}
                    yield f"data: {json.dumps(page_data)}\n\n"
                    yield f"data: {json.dumps({'done': True, 'note_name': name})}\n\n"
                
                elif file_type == "txt" or file_type == "md":
                    if is_encoded:
                        text = decoded_content.decode('utf-8', errors='ignore')
                    else:
                        text = content
                    text = clean_extracted_text(text)
                    
                    page_data = {"page_num": 1, "total_pages": 1, "text": text, "extracted": True}
                    yield f"data: {json.dumps(page_data)}\n\n"
                    yield f"data: {json.dumps({'done': True, 'note_name': name})}\n\n"
                
                elif file_type in ["ppt", "pptx"]:
                    try:
                        from pptx import Presentation
                        doc_bytes = decoded_content if is_encoded else content.encode()
                        doc_file = BytesIO(doc_bytes)
                        prs = Presentation(doc_file)
                        total_slides = len(prs.slides)
                        
                        for slide_num, slide in enumerate(prs.slides, 1):
                            slide_text = f"\n=== Slide {slide_num} ===\n\n"
                            for shape in slide.shapes:
                                if hasattr(shape, "text_frame"):
                                    for para in shape.text_frame.paragraphs:
                                        para_text = para.text.strip()
                                        if para_text:
                                            slide_text += para_text + "\n"
                                elif hasattr(shape, "text") and shape.text:
                                    slide_text += shape.text.strip() + "\n"
                            if slide_text.strip():
                                text = clean_extracted_text(slide_text)
                                page_data = {
                                    "page_num": slide_num,
                                    "total_pages": total_slides,
                                    "text": text,
                                    "extracted": True
                                }
                                yield f"data: {json.dumps(page_data)}\n\n"
                                await asyncio.sleep(0.05)
                        
                        yield f"data: {json.dumps({'done': True, 'note_name': name})}\n\n"
                    except Exception as e:
                        yield f"data: {json.dumps({'error': str(e)})}\n\n"
                
                else:
                    yield f"data: {json.dumps({'error': 'Unsupported format'})}\n\n"
                
                yield f"data: {json.dumps({'done': True})}\n\n"
            
            except Exception as e:
                yield f"data: {json.dumps({'error': str(e)})}\n\n"
        
        return StreamingResponse(generate_extraction(), media_type="text/event-stream")
    
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/notes")
async def save_note(
    user: str = Form(...),
    name: str = Form(...),
    content: str = Form(...),
    file_type: str = Form(...)
):
    try:
        note_id = save_lecture_note(user, name, content, file_type)
        return {"success": True, "note_id": note_id}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.delete("/api/notes/{note_id}")
async def delete_note(note_id: int):
    try:
        delete_lecture_note(note_id)
        return {"success": True}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

# LangChain-powered endpoints
@app.post("/api/chat/generate-quiz")
async def generate_quiz(
    topic: str = Form(...),
    difficulty: str = Form(default="Medium"),
    num_questions: int = Form(default=5),
    user: str = Form(default="User"),
    use_rag: bool = Form(default=False)
):
    try:
        if use_rag:
            result = await generate_quiz_with_rag(topic, difficulty, num_questions, user, 0)
        else:
            from api.services.langchain_service import get_langchain_llm
            from langchain_core.prompts import ChatPromptTemplate
            
            llm = get_langchain_llm()
            prompt = ChatPromptTemplate.from_messages([
                ("system", f"You are a quiz generator. Create a {difficulty} quiz with {num_questions} multiple choice questions about the given topic.\nFormat each question as:\nQ1) [question]\nA) [option1]\nB) [option2]\nC) [option3]\nD) [option4]\nAnswer: [correct answer letter]"),
                ("human", "Topic: {topic}")
            ])
            chain = prompt | llm
            response = chain.invoke({"topic": topic})
            result = response.content if hasattr(response, 'content') else str(response)
        
        return {"quiz": result}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/chat/generate-flashcards")
async def generate_flashcards(
    topic: str = Form(...),
    num_cards: int = Form(default=5),
    user: str = Form(default="User"),
    use_rag: bool = Form(default=False)
):
    try:
        if use_rag:
            result = await generate_flashcards_with_rag(topic, num_cards, user, 0)
        else:
            from api.services.langchain_service import get_langchain_llm
            from langchain_core.prompts import ChatPromptTemplate
            
            llm = get_langchain_llm()
            prompt = ChatPromptTemplate.from_messages([
                ("system", f"You are an expert educator. Create {num_cards} flashcards about the given topic.\nFormat each card as:\nQ: [question]\nA: [answer]"),
                ("human", "Topic: {topic}")
            ])
            chain = prompt | llm
            response = chain.invoke({"topic": topic})
            result = response.content if hasattr(response, 'content') else str(response)
        
        return {"flashcards": result}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/rag/ingest")
async def ingest_document(
    file: UploadFile = File(...),
    doc_name: str = Form(default="document")
):
    try:
        content = await file.read()
        content_type = file.content_type or ""
        
        import base64
        b64_data = base64.b64encode(content).decode("utf-8")
        
        if content_type.startswith("image/"):
            return JSONResponse({"error": "Image RAG not supported yet"}, status_code=400)
        
        from api.services.langchain_service import rag_system
        
        if file.filename.endswith(".pdf"):
            pdf_bytes = base64.b64decode(b64_data)
            success = get_rag_system().process_pdf(pdf_bytes, doc_name)
        else:
            text = content.decode("utf-8", errors="ignore")
            success = get_rag_system().process_text(text, doc_name)
        
        if success:
            return {"success": True, "document": doc_name, "chunks": "created"}
        return JSONResponse({"error": "Failed to create embeddings"}, status_code=500)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/api/rag/search")
async def search_documents(
    query: str,
    doc_name: Optional[str] = None,
    k: int = 3
):
    try:
        results = get_rag_system().similarity_search(query, doc_name, k)
        return {"results": results}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)